from langchain_community.document_loaders import PyMuPDFLoader
import docx
from pptx import Presentation


def extract_text(path):
    """ 没有使用 utf-8 编码打开，返回的是二进制数据 """
    return open(path, 'r').read()


def extract_pdf(path) -> str:
    """ 使用 langchain 提供的 PyMuPDF 工具按页读取 pdf 文件 """
    loader = PyMuPDFLoader(path)
    data = loader.load()
    data = [x.page_content for x in data]
    content = '\n\n'.join(data)
    return content


def extract_docx(path) -> str:
    """ 使用 python-docx 三方件逐段落读取 doc 文档内容 """
    doc = docx.Document(path)
    data = []
    for paragraph in doc.paragraphs:
        data.append(paragraph.text)
    content = '\n\n'.join(data)
    return content


def extract_pptx(path) -> str:
    """ 使用 python-docx 三方件逐页读取 ppt 的文本内容 """
    prs = Presentation(path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text
